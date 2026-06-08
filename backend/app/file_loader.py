from pathlib import Path
from tempfile import NamedTemporaryFile

from docx import Document
from fastapi import UploadFile
from pypdf import PdfReader


SUPPORTED_EXTENSIONS = {".txt", ".md", ".pdf", ".docx", ".csv", ".ppt", ".pptx"}


async def extract_text(upload: UploadFile) -> str:
    suffix = Path(upload.filename or "").suffix.lower()
    raw = await upload.read()

    with NamedTemporaryFile(delete=False, suffix=suffix) as temp_file:
        temp_file.write(raw)
        temp_path = Path(temp_file.name)

    try:
        if suffix == ".pdf":
            reader = PdfReader(str(temp_path))
            return "\n".join(page.extract_text() or "" for page in reader.pages)

        if suffix == ".docx":
            document = Document(str(temp_path))
            return "\n".join(paragraph.text for paragraph in document.paragraphs)

        if suffix == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(temp_path))
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_runs.append(shape.text)
            return "\n".join(text_runs)

        if suffix == ".ppt":
            raise ValueError(
                "Legacy .ppt format is not directly supported. "
                "Please save the file as a modern .pptx presentation and try uploading again."
            )

        # Fallback for all other files (txt, md, csv, code files, unknown)
        return raw.decode("utf-8", errors="ignore")
    finally:
        temp_path.unlink(missing_ok=True)

